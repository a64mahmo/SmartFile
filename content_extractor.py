from pathlib import Path
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import logging
from typing import Optional
import subprocess
import tempfile
import os
import shutil
from AppKit import NSWorkspace
import pandas as pd
import zipfile
import xml.etree.ElementTree as ET
import re

class ContentExtractor:
    def __init__(self):
        self.logger = logging.getLogger(__name__)
        self.max_content_length = 1000000  # Increased limit for better content extraction
        self.temp_dir = Path(tempfile.mkdtemp())

    def __del__(self):
        """Clean up temporary directory when the object is destroyed."""
        try:
            shutil.rmtree(self.temp_dir)
        except Exception as e:
            self.logger.error(f"Error cleaning up temp directory: {str(e)}")

    def extract_content(self, file_path: Path) -> Optional[str]:
        """Extract text content from a file based on its extension."""
        try:
            if not file_path.exists():
                self.logger.error(f"File does not exist: {file_path}")
                return None

            self.logger.info(f"Extracting content from: {file_path}")
            content = ""
            
            if file_path.suffix.lower() == '.pdf':
                content = self._extract_pdf_content(file_path)
            elif file_path.suffix.lower() == '.docx':
                content = self._extract_docx_content(file_path)
            elif file_path.suffix.lower() == '.txt':
                content = self._extract_text_content(file_path)
            elif file_path.suffix.lower() == '.pptx':
                content = self._extract_pptx_content(file_path)
            elif file_path.suffix.lower() == '.dmg':
                content = self._extract_dmg_content(file_path)
            elif file_path.suffix.lower() == '.xlsx':
                content = self._extract_xlsx_content(file_path)
            elif file_path.suffix.lower() == '.epub':
                content = self._extract_epub_content(file_path)
            else:
                self.logger.warning(f"Unsupported file type: {file_path.suffix}")
                return None

            if not content:
                self.logger.warning(f"No content extracted from {file_path}")
                return None

            # Truncate content if it exceeds the maximum length
            if len(content) > self.max_content_length:
                self.logger.warning(f"Content truncated from {len(content)} to {self.max_content_length} characters")
                content = content[:self.max_content_length]

            return content

        except Exception as e:
            self.logger.error(f"Error extracting content from {file_path}: {str(e)}")
            return None

    def _extract_pdf_content(self, file_path: Path) -> Optional[str]:
        """Extract text content from a PDF file using PyMuPDF."""
        try:
            doc = fitz.open(file_path)
            text = ""
            
            for page in doc:
                text += page.get_text()
            
            doc.close()
            return text if text.strip() else None
            
        except Exception as e:
            self.logger.error(f"Error extracting PDF content from {file_path}: {str(e)}")
            return None

    def _extract_docx_content(self, file_path: Path) -> Optional[str]:
        """Extract text content from a DOCX file."""
        try:
            doc = docx.Document(file_path)
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
        except Exception as e:
            self.logger.error(f"Error extracting DOCX content from {file_path}: {str(e)}")
            return None

    def _extract_text_content(self, file_path: Path) -> Optional[str]:
        """Extract content from a text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='latin-1') as f:
                    return f.read()
            except Exception as e:
                self.logger.error(f"Error extracting text content from {file_path}: {str(e)}")
                return None
        except Exception as e:
            self.logger.error(f"Error extracting text content from {file_path}: {str(e)}")
            return None

    def _extract_pptx_content(self, file_path: Path) -> Optional[str]:
        """Extract text content from a PPTX file."""
        try:
            prs = Presentation(file_path)
            text = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text.append(paragraph.text)
            
            return "\n".join(text) if text else None
            
        except Exception as e:
            self.logger.error(f"Error extracting PPTX content from {file_path}: {str(e)}")
            return None

    def _extract_dmg_content(self, file_path: Path) -> Optional[str]:
        """Extract basic information from a DMG file."""
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                # Mount the DMG
                mount_cmd = ["hdiutil", "attach", "-nobrowse", "-mountpoint", temp_dir, str(file_path)]
                subprocess.run(mount_cmd, check=True)
                
                # Get basic information
                info = []
                for root, dirs, files in os.walk(temp_dir):
                    for file in files:
                        info.append(f"File: {file}")
                    for dir in dirs:
                        info.append(f"Directory: {dir}")
                
                # Unmount the DMG
                unmount_cmd = ["hdiutil", "detach", temp_dir]
                subprocess.run(unmount_cmd, check=True)
                
                return "\n".join(info) if info else None
                
        except Exception as e:
            self.logger.error(f"Error extracting DMG content from {file_path}: {str(e)}")
            return None

    def _extract_xlsx_content(self, file_path: Path) -> Optional[str]:
        """Extract text content from an XLSX file."""
        try:
            # Read all sheets from the Excel file
            excel_file = pd.ExcelFile(file_path)
            content = []
            
            # Process each sheet
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # Add sheet name as header
                content.append(f"\nSheet: {sheet_name}\n")
                
                # Add column headers
                content.append(" | ".join(str(col) for col in df.columns))
                content.append("-" * 50)
                
                # Add data rows
                for _, row in df.iterrows():
                    content.append(" | ".join(str(val) for val in row))
            
            return "\n".join(content) if content else None
            
        except Exception as e:
            self.logger.error(f"Error extracting XLSX content from {file_path}: {str(e)}")
            return None

    def _extract_epub_content(self, file_path: Path) -> Optional[str]:
        """Extract text content from an EPUB file."""
        try:
            content = []
            
            # Open the EPUB file
            with zipfile.ZipFile(file_path, 'r') as epub:
                # Get the container file to find the content files
                container = ET.fromstring(epub.read('META-INF/container.xml'))
                rootfile = container.find('.//{urn:oasis:names:tc:opendocument:xmlns:container}rootfile')
                if rootfile is None:
                    self.logger.error("Could not find rootfile in EPUB container")
                    return None
                    
                # Get the content file path
                content_path = rootfile.get('full-path')
                if not content_path:
                    self.logger.error("Could not find content path in EPUB")
                    return None
                    
                # Read the content file
                content_xml = ET.fromstring(epub.read(content_path))
                
                # Extract text from all text elements
                for elem in content_xml.iter():
                    if elem.text and elem.text.strip():
                        content.append(elem.text.strip())
                    
            return '\n'.join(content) if content else None
            
        except Exception as e:
            self.logger.error(f"Error extracting EPUB content from {file_path}: {str(e)}")
            return None 